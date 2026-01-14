from pptx import Presentation
import re
from datetime import datetime
from pathlib import Path
import pandas as pd


BASE_DIR = Path(__file__).resolve().parent # report folder
GUI_DIR = BASE_DIR.parent  # GUI folder
RESULTS_ROOT = GUI_DIR.parent / "battery_optimization" / "results"
TEMPLATE_FOLDER = GUI_DIR / "templates" 


company = "039. Godestadsvagen"
correspondance_name = "Gödestadsvägen"
date_folder = "2026-01-12"
run_id = "5_7f53f1fe-d366-47eb-882a-fbd0c66c0296"

RESULTS_PATH = RESULTS_ROOT / company / date_folder / run_id / "Results_All.xlsx"

# select template choose what you want to show to customer by eying the results excel and selecting the appropriate template
template_key = "ownPVss"  # choose from: ownPVss, ownPV, APIPVss, APIPV          ,ss = self-sufficient 


df_results = pd.read_excel(RESULTS_PATH)
print(df_results.head(19))

# KEEP IN MIND 0-BASED INDEXING FOR ROW SELECTION 
# todo: run baseline cases 

row_selected_case_2024 = 1  # not self-sufficient, 2024 (_1)
row_selected_case_2024_ss = 5  # self-sufficient, 2024 (_2)
row_selected_case_2025 = 9  # not self-sufficient, 2025 (_3)
row_selected_case_2025_ss = 13  # self-sufficient, 2025 (_4)

templates = {
            "ownPVss": TEMPLATE_FOLDER / "template_ownPVss.pptx",
            "ownPV": TEMPLATE_FOLDER / "template_ownPV.pptx",
            "APIPVss": TEMPLATE_FOLDER / "template_APIPVss.pptx",
            "APIPV": TEMPLATE_FOLDER / "template_APIPV.pptx",
}


# add in format: additional_information_<correspondance_name>
additional_information = {
    "Poseidon": "Net hourly demand data was provided by Poseidon together with a PV specifications sheet and the PV area. "
    "Based on this information, we estimated 205 kWp of installed solar power. Furthermore, gross hourly demand "
    "data was calculated with the help of hourly means. As importance of sustainability objectives was expressed, "
    "this version includes a column where the self-sufficiency of the system is prioritised (i.e. % of total own "
    "electricity demand met by PV generation).",

    "Herrestad": "Information",

    "Gödestadsvägen": "Net hourly demand data was provided together with hourly PV production for 2024 and 2025. "
    "Based on this information, gross hourly demand data was calculated with the help of hourly means. "
}

PATTERN = re.compile(r"\{\{\s*([\w.]+)\s*\}\}")

def readable_demand_timestamps(s: str):
    start_str, end_str = s.split(" to ")
    start = datetime.strptime(start_str, "%Y-%m-%d %H:%M:%S")
    end = datetime.strptime(end_str, "%Y-%m-%d %H:%M:%S")
    return f"{start.strftime('%b %Y')} - {end.strftime('%b %Y')}"

def readable_demand_timestamps_year(s:str):
    end = datetime.strptime(s, "%Y-%m-%d %H:%M:%S")
    return f"{end.strftime('%Y')}"

def readable_price_timestamps(s: str, e: str):
    start = datetime.strptime(s, "%Y-%m-%d %H:%M:%S")
    end = datetime.strptime(e, "%Y-%m-%d %H:%M:%S")
    return f"{start.strftime('%b %Y')} - {end.strftime('%b %Y')}"

def replace_placeholders_pptx(template_pptx, output_pptx, context):
    prs = Presentation(template_pptx)

    def replace_in_paragraph(paragraph):
        if not paragraph.runs:
            return

        full = "".join(r.text for r in paragraph.runs)
        if "{{" not in full:
            return

        # remove invisible chars that can break matching
        full_clean = full.replace("\u200b", "").replace("\ufeff", "")

        def repl(m):
            key = m.group(1).strip()
            return str(context.get(key, m.group(0)))

        new = PATTERN.sub(repl, full_clean)
        if new == full_clean:
            return

        # Write result into first run, clear rest
        paragraph.runs[0].text = new
        for r in paragraph.runs[1:]:
            r.text = ""

    def replace_in_textframe(tf):
        for p in tf.paragraphs:
            replace_in_paragraph(p)

    def walk_shapes(shapes):
        for shape in shapes:
            # Grouped shapes
            if shape.shape_type == 6:  # GROUP
                walk_shapes(shape.shapes)
                continue

            # Text frames
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                replace_in_textframe(shape.text_frame)

            # Tables
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_in_textframe(cell.text_frame)

    # Slides
    for slide in prs.slides:
        walk_shapes(slide.shapes)

    # Layouts (important)
    for layout in prs.slide_layouts:
        walk_shapes(layout.shapes)

    # Masters (important)
    for master in prs.slide_masters:
        walk_shapes(master.shapes)

    prs.save(output_pptx)

def build_case_context(row, correspondance_name: str, additional_information: dict) -> dict:
    return {
        "data_year": readable_demand_timestamps_year(row["end_time"]),
        "demand_data_range": readable_demand_timestamps(row["demand_data_used:"]),
        "price_data_range": readable_price_timestamps(row["start_time"], row["end_time"]),
        "simulation_length": round(row["length of simulation"], 3),

        "cap": row["Battery Capacity"],
        "pow": row["Max Power"],

        "c_b": round(row["Total electricity cost base case"] / row["length of simulation"], 0),
        "c_pv": round(row["Total electricity cost with PV"] / row["length of simulation"], 0),
        "c_pv_bat": round(row["Total electricity cost with PV and battery"] /row["length of simulation"], 0),

        "c_k_b": round(row["Electricity cost/kWh base case"], 3),
        "c_k_pv": round(row["Electricity cost/kWh with PV"], 3),
        "c_k_pvb": round(row["Electricity cost/kWh with PV and battery"], 4),

        "d_p": round(row["€ % drop/kWh"], 1),  # drop in €/kWh
        "d_pv": round((1-row["Electricity cost/kWh with PV and battery"]/row["Electricity cost/kWh with PV"])*100, 1),  # drop in €/kWh for own PV


        # savings
        "s_pv": round(row["Savings from PV during period"] / row["length of simulation"], 0),
        "s_b": round(row["Savings from battery during period"]/ row["length of simulation"], 0),
        "s_t": round(row["Savings from (PV and) battery during period"]/ row["length of simulation"], 0),
        "s_pcf": round(row["Savings from reduced power contract fee"]/ row["length of simulation"], 0),

        "pb_pv": round(row["Payback Time PV"], 1),
        "pb_b": round(row["Payback Time battery"], 1),
        "pb_t": round(row["Payback Time (PV and) battery"], 1),

        "ssr_pv": round(row["Self-sufficiency (SSR only PV %)"], 1),
        "ssr_t": round(row["Self-sufficiency (SSR %)"], 1),  
        "self_consumption": round(row["Self-consumption (%)"], 1),

        "cycles": round(row["Number of Cycles in period"] / row["length of simulation"], 1),
        "time_below_20": round(row["Time below 20% SoC (%)"], 1),
        "time_above_80": round(row["Time above 80% SoC (%)"], 1),
        "t_20_80": round(row["%time_SoC_20%-80%"], 1),
        "charge": round((row["Number of Cycles in period"] / (1000*row["length of simulation"])*row["Battery Capacity"]), 1),

        "b_in": round(row["Battery investment (EUR)"], 0),

        "pv_price": round(row["PV Price (EUR/kW)"], 0),
        "bat_energy_price": round(row["Battery Capacity Price (EUR/kWh)"], 0),
        "bat_power_price": round(row["Battery Power Price (EUR/kWh)"], 0),

        "fcr_rev": round(row["Total FCR revenue (EUR)"]/row["length of simulation"], 0),
        
    }

def suffix_keys(d: dict, suffix: str) -> dict:
    return {f"{k}{suffix}": v for k, v in d.items()}



if __name__ == "__main__":
    case_indices = [
        row_selected_case_2024,
        row_selected_case_2024_ss,
        row_selected_case_2025,
        row_selected_case_2025_ss,
    ]
    # Metrics
    
    pps = round(
        sum(df_results.iloc[i]["peak% of saving"] for i in case_indices) / len(case_indices), 1 )
    pamr = round(
        sum(df_results.iloc[i]["FCR% of saving"] for i in case_indices) / len(case_indices),  1)
    parb = round(100 - pps - pamr, 1)

    # Use the first selected case as the "base" row for shared values
    base_row = df_results.iloc[case_indices[0]]

    base_context = {
        "additional_information": additional_information.get(correspondance_name, ""),
        "cor_na": correspondance_name,
        "region": base_row["Region"],
        "data_year": readable_demand_timestamps_year(base_row["end_time"]),
        "demand_data_range": readable_demand_timestamps(base_row["demand_data_used:"]),
        "price_data_range": readable_price_timestamps(
            base_row["start_time"], base_row["end_time"]
        ),
        "pv_size_kw": (df_results.iloc[row_selected_case_2024]["PV Size"]),                 # installed PV
        "efficiency": round(base_row["Energy Efficiency"] * 100, 1),
        "pv_in": round(base_row["PV investment (EUR)"], 0),

    #   2024 and 2025 average values
        "total_demand": round((df_results.iloc[row_selected_case_2024]["Total demand (kWh)"] + df_results.iloc[row_selected_case_2025]["Total demand (kWh)"] )/ (2 * 1000),1),
        "solar_mwh": round((df_results.iloc[row_selected_case_2024]["Total PV production during period"] + df_results.iloc[row_selected_case_2025]["Total PV production during period"] )/ (2 ),1),
        "PPS" :pps,
        "PARB": parb,
        "PAMR": pamr,
        "ss_incr2":  round((df_results.iloc[row_selected_case_2024_ss]["Self-sufficiency (SSR %)"])/df_results.iloc[row_selected_case_2024]["Self-sufficiency (SSR only PV %)"],1),
        "ss_incr":  round((df_results.iloc[row_selected_case_2025_ss]["Self-sufficiency (SSR %)"])/df_results.iloc[row_selected_case_2025]["Self-sufficiency (SSR only PV %)"],1),

    }


    context_all = base_context.copy()

    for i, idx in enumerate(case_indices, start=1):
        row = df_results.iloc[idx]
        case_ctx = build_case_context(row, correspondance_name, additional_information)
        context_all.update(suffix_keys(case_ctx, f"_{i}"))


    replace_placeholders_pptx(
        template_pptx = templates[template_key],
        output_pptx=fr"C:\Users\meiek\Documents\RIVUS_code\GUI\report\output_{correspondance_name}.pptx",
        context=context_all,
    )
